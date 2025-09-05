from fastapi import APIRouter, UploadFile, File, HTTPException, status
from typing import List
import fitz  # PyMuPDF for PDF
import pytesseract
from PIL import Image
import io
import docx
import pptx
import tempfile
import os

router = APIRouter()

SUPPORTED_TYPES = [
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # Word
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # PowerPoint
    "text/plain",
    "image/png",
    "image/jpeg"
]
MAX_FILE_SIZE_MB = 10

def extract_pdf_text(file_bytes: bytes) -> str:
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        if text.strip():  # If text was extracted
            return text
        # If PDF is scanned/image, use OCR
        images = []
        for page in doc:
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            images.append(img)
        ocr_text = ""
        for img in images:
            ocr_text += pytesseract.image_to_string(img)
        return ocr_text
    except Exception as e:
        raise RuntimeError(f"PDF extraction failed: {e}")

def extract_word_text(file_bytes: bytes) -> str:
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = docx.Document(tmp_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        os.remove(tmp_path)
        return text
    except Exception as e:
        raise RuntimeError(f"Word extraction failed: {e}")

def extract_ppt_text(file_bytes: bytes) -> str:
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        prs = pptx.Presentation(tmp_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        os.remove(tmp_path)
        return text
    except Exception as e:
        raise RuntimeError(f"PPT extraction failed: {e}")

def extract_image_text(file_bytes: bytes) -> str:
    try:
        img = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        raise RuntimeError(f"Image OCR failed: {e}")

def extract_plain_text(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode(errors="ignore")
    except Exception as e:
        raise RuntimeError(f"Plain text extraction failed: {e}")

@router.post("/upload-material", status_code=status.HTTP_201_CREATED)
async def upload_material(file: UploadFile = File(...)) -> dict:
    """
    Accepts PDF, Word, PowerPoint, plain text, or image files.
    Extracts text from the file and returns it for user review.
    """
    if file.content_type not in SUPPORTED_TYPES:
        raise HTTPException(status_code=415, detail=f"Unsupported file type: {file.content_type}")

    contents = await file.read()
    file_size_mb = len(contents) / (1024 * 1024)
    if file_size_mb > MAX_FILE_SIZE_MB:
        raise HTTPException(status_code=413, detail=f"File too large. Max allowed is {MAX_FILE_SIZE_MB} MB.")

    try:
        if file.content_type == "application/pdf":
            text = extract_pdf_text(contents)
        elif file.content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text = extract_word_text(contents)
        elif file.content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = extract_ppt_text(contents)
        elif file.content_type.startswith("image/"):
            text = extract_image_text(contents)
        elif file.content_type == "text/plain":
            text = extract_plain_text(contents)
        else:
            raise HTTPException(status_code=415, detail="File type not supported.")
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Text extraction failed: {e}")

    if not text.strip():
        raise HTTPException(status_code=422, detail="No text could be extracted from file. Try another file or check file quality.")

    return {
        "filename": file.filename,
        "filetype": file.content_type,
        "filesize_mb": round(file_size_mb, 2),
        "extracted_text": text[:5000],  # Return only first 5000 chars for preview
        "full_text_available": len(text) > 5000
    }